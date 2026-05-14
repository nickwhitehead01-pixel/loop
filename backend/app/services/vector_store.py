"""
Document ingestion and ChromaDB similarity search.

ingest_lesson  — parse a document, chunk with 500/50 token splits, embed, store LessonChunks
search         — embed a query, return top-k chunk texts
"""
from __future__ import annotations

import asyncio
import base64
import logging
from io import BytesIO

from langchain_text_splitters import RecursiveCharacterTextSplitter
from pypdf import PdfReader
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.domain import LessonChunk
from app.services.chroma_client import lesson_chunks_col
from app.services.ollama_client import describe_image, embed, embed_batch

logger = logging.getLogger(__name__)

# Chunks smaller than this (chars) are filtered after splitting — avoids noise
MIN_CHUNK_CHARS = 80

# 500-token chunks with 50-token overlap using the GPT-2 tokenizer
# (closest public approximation to nomic-embed-text's vocabulary size)
_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
    encoding_name="gpt2",
    chunk_size=500,
    chunk_overlap=50,
)

ACCEPTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


# ---------------------------------------------------------------------------
# Text sanitization — remove invalid UTF-8 characters
# ---------------------------------------------------------------------------

def _sanitize_text(text: str) -> str:
    """Remove null bytes and other problematic characters that break UTF-8 encoding."""
    return "".join(char if ord(char) >= 32 or char in "\n\r\t" else " " for char in text)


# ---------------------------------------------------------------------------
# Text extraction — one function per format
# ---------------------------------------------------------------------------

def _extract_pdf(data: bytes) -> str:
    reader = PdfReader(BytesIO(data))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            text = _sanitize_text(text)
            pages.append(text)
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    from docx import Document  # python-docx
    doc = Document(BytesIO(data))
    paragraphs = [_sanitize_text(p.text) for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation  # python-pptx
    prs = Presentation(BytesIO(data))
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = _sanitize_text(para.text.strip())
                    if line:
                        texts.append(line)
        if texts:
            slides.append("\n".join(texts))
    return "\n\n".join(slides)


def _extract_txt(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    return _sanitize_text(text)


# ---------------------------------------------------------------------------
# Image extraction — one function per format
# Returns list of (location, image_bytes) where location is a 1-based
# page/slide number (0 when location is unavailable, e.g. DOCX).
# Images smaller than 2 KB are skipped (icons, decorative bullets, etc.).
# ---------------------------------------------------------------------------

_MIN_IMAGE_BYTES = 2048  # skip anything smaller than ~2 KB


def _extract_images_pdf(data: bytes) -> list[tuple[int, bytes]]:
    """Extract unique images from each PDF page via PyMuPDF."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("pymupdf not installed — PDF image extraction skipped")
        return []
    images: list[tuple[int, bytes]] = []
    pdf = fitz.open(stream=data, filetype="pdf")
    try:
        seen_xrefs: set[int] = set()
        for page_idx, page in enumerate(pdf):
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                try:
                    base_image = pdf.extract_image(xref)
                    img_bytes = base_image["image"]
                    if len(img_bytes) >= _MIN_IMAGE_BYTES:
                        images.append((page_idx + 1, img_bytes))
                except Exception:
                    logger.debug("PDF image extraction failed for page %d xref %d — skipping", page_idx + 1, xref, exc_info=True)
                    continue
    finally:
        pdf.close()
    return images


def _extract_images_pptx(data: bytes) -> list[tuple[int, bytes]]:
    """Extract picture shapes from each PPTX slide."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    images: list[tuple[int, bytes]] = []
    prs = Presentation(BytesIO(data))
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    if len(img_bytes) >= _MIN_IMAGE_BYTES:
                        images.append((slide_idx + 1, img_bytes))
                except Exception:
                    logger.debug("PPTX image extraction failed for slide %d — skipping", slide_idx + 1, exc_info=True)
                    continue
    return images


def _extract_images_docx(data: bytes) -> list[tuple[int, bytes]]:
    """Extract inline images from a DOCX document."""
    from docx import Document

    images: list[tuple[int, bytes]] = []
    doc = Document(BytesIO(data))
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_bytes = rel.target_part.blob
                if len(img_bytes) >= _MIN_IMAGE_BYTES:
                    images.append((0, img_bytes))
            except Exception:
                logger.debug("DOCX image extraction failed — skipping", exc_info=True)
                continue
    return images


def _extract_images(file_bytes: bytes, filename: str) -> list[tuple[int, bytes]]:
    """Dispatch to the correct image extractor based on file extension."""
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""
    if ext == ".pdf":
        return _extract_images_pdf(file_bytes)
    if ext == ".pptx":
        return _extract_images_pptx(file_bytes)
    if ext == ".docx":
        return _extract_images_docx(file_bytes)
    return []


def extract_text(file_bytes: bytes, filename: str) -> str:
    """
    Dispatch to the correct parser based on file extension.
    Raises ValueError for unsupported types.
    """
    ext = ""
    if "." in filename:
        ext = "." + filename.rsplit(".", 1)[-1].lower()

    if ext == ".pdf":
        return _extract_pdf(file_bytes)
    elif ext == ".docx":
        return _extract_docx(file_bytes)
    elif ext == ".pptx":
        return _extract_pptx(file_bytes)
    elif ext == ".txt":
        return _extract_txt(file_bytes)
    else:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Accepted: {', '.join(sorted(ACCEPTED_EXTENSIONS))}"
        )


def _chunk_text(text: str) -> list[str]:
    """
    Split *text* into 500-token chunks with 50-token overlap.
    Uses RecursiveCharacterTextSplitter with the GPT-2 tiktoken encoder.
    Chunks shorter than MIN_CHUNK_CHARS chars are discarded (page numbers, etc.).
    """
    raw_chunks = _splitter.split_text(text)
    return [c for c in raw_chunks if len(c.strip()) >= MIN_CHUNK_CHARS]


# ---------------------------------------------------------------------------
# Located text extraction — returns (chunk_text, slide_number) pairs so that
# every chunk carries a 1-based page/slide index for semantic slide-sync.
# slide_number = 0 means "position unknown" (DOCX, TXT).
# ---------------------------------------------------------------------------

def _extract_pdf_located(data: bytes) -> list[tuple[str, int]]:
    """Return one (page_text, page_number) entry per PDF page."""
    reader = PdfReader(BytesIO(data))
    result: list[tuple[str, int]] = []
    for page_idx, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            result.append((_sanitize_text(text), page_idx + 1))
    return result


def _extract_pptx_located(data: bytes) -> list[tuple[str, int]]:
    """Return one (slide_text, slide_number) entry per PPTX slide."""
    from pptx import Presentation  # python-pptx

    prs = Presentation(BytesIO(data))
    result: list[tuple[str, int]] = []
    for slide_idx, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = _sanitize_text(para.text.strip())
                    if line:
                        texts.append(line)
        if texts:
            result.append(("\n".join(texts), slide_idx + 1))
    return result


def _extract_located(file_bytes: bytes, filename: str) -> list[tuple[str, int]]:
    """Dispatch to a per-page/per-slide extractor based on file extension.

    DOCX and TXT have no page/slide structure so the entire text is returned
    as a single entry with slide_number = 0 ("position unknown").
    """
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""
    if ext == ".pptx":
        return _extract_pptx_located(file_bytes)
    if ext == ".pdf":
        return _extract_pdf_located(file_bytes)
    # DOCX / TXT — no positional structure
    text = extract_text(file_bytes, filename)
    return [(text, 0)] if text.strip() else []


def _chunk_text_located(texts_with_locations: list[tuple[str, int]]) -> list[tuple[str, int]]:
    """Chunk each (text, location) independently, propagating the location to every chunk."""
    result: list[tuple[str, int]] = []
    for text, loc in texts_with_locations:
        for chunk in _splitter.split_text(text):
            if len(chunk.strip()) >= MIN_CHUNK_CHARS:
                result.append((chunk, loc))
    return result


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

async def ingest_lesson(
    lesson_id: int,
    file_bytes: bytes,
    filename: str = "upload.pdf",
) -> int:
    """
    Parse *file_bytes*, embed all text chunks, and store into lesson_chunks.
    Returns the number of text chunks stored.

    Manages its own DB session so it can be called as a background task from
    the upload endpoint, letting the HTTP response return immediately after
    the Lesson record is created (the UI's "Reading file…" stage covers this).

    Text extraction (pypdf/pptx) and ChromaDB writes are run in a thread
    pool so they don't block the async event loop.
    """
    from app.core.database import AsyncSessionLocal

    # CPU-bound extraction runs off the event loop
    located_text_chunks = await asyncio.to_thread(
        lambda: _chunk_text_located(_extract_located(file_bytes, filename))
    )
    if not located_text_chunks:
        return 0

    text_chunks, text_slide_numbers = zip(*located_text_chunks)
    text_chunks = list(text_chunks)
    text_slide_numbers = list(text_slide_numbers)

    vectors = await embed_batch(text_chunks)

    async with AsyncSessionLocal() as db:
        for i, chunk_text in enumerate(text_chunks):
            db.add(LessonChunk(lesson_id=lesson_id, chunk_index=i, content=chunk_text))

        await db.flush()

        result = await db.execute(
            select(LessonChunk.id, LessonChunk.content, LessonChunk.chunk_index)
            .where(LessonChunk.lesson_id == lesson_id)
            .order_by(LessonChunk.chunk_index)
        )
        rows = result.all()

        col = lesson_chunks_col()
        chroma_ids: list[str] = []
        chroma_embeddings: list[list[float]] = []
        chroma_documents: list[str] = []
        chroma_metadatas: list[dict] = []

        for (chunk_id, chunk_content, chunk_index), vector in zip(rows, vectors):
            chroma_ids.append(str(chunk_id))
            chroma_embeddings.append(vector)
            chroma_documents.append(chunk_content)
            chroma_metadatas.append({
                "lesson_id": str(lesson_id),
                "chunk_index": chunk_index,
                "slide_number": str(text_slide_numbers[chunk_index]),
            })

        # Sync ChromaDB write — run in thread pool to avoid blocking the loop
        await asyncio.to_thread(
            lambda: col.add(
                ids=chroma_ids,
                embeddings=chroma_embeddings,
                documents=chroma_documents,
                metadatas=chroma_metadatas,
            )
        )

        await db.commit()

    logger.info("lesson_id=%s ingested: %d text chunk(s)", lesson_id, len(rows))
    return len(rows)


async def ingest_lesson_images(
    lesson_id: int,
    file_bytes: bytes,
    filename: str = "upload.pdf",
    chunk_offset: int = 0,
) -> int:
    """
    Describe images extracted from *file_bytes* using gemma4 vision and add
    the descriptions as additional chunks in ChromaDB.

    Designed to be called as a background task AFTER ingest_lesson() returns,
    so the upload endpoint stays fast. chunk_offset should be the number of
    text chunks already stored (returned by ingest_lesson) so indices don't
    collide. Uses its own DB session.
    """
    from app.core.database import AsyncSessionLocal

    raw_images = _extract_images(file_bytes, filename)
    if not raw_images:
        return 0

    image_chunks: list[str] = []
    image_slide_numbers: list[int] = []

    for location, img_bytes in raw_images:
        try:
            img_b64 = base64.b64encode(img_bytes).decode("ascii")
            description = await describe_image(img_b64)
            if description:
                image_chunks.append(description)
                image_slide_numbers.append(location)
        except Exception:
            logger.warning(
                "lesson_id=%s — image description failed (skipping)", lesson_id, exc_info=True
            )

    if not image_chunks:
        return 0

    vectors = await embed_batch(image_chunks)

    async with AsyncSessionLocal() as db:
        for i, chunk_text in enumerate(image_chunks):
            db.add(LessonChunk(
                lesson_id=lesson_id,
                chunk_index=chunk_offset + i,
                content=chunk_text,
            ))
        await db.flush()

        result = await db.execute(
            select(LessonChunk.id, LessonChunk.content, LessonChunk.chunk_index)
            .where(LessonChunk.lesson_id == lesson_id)
            .where(LessonChunk.chunk_index >= chunk_offset)
            .order_by(LessonChunk.chunk_index)
        )
        rows = result.all()

        col = lesson_chunks_col()
        chroma_ids: list[str] = []
        chroma_embeddings: list[list[float]] = []
        chroma_documents: list[str] = []
        chroma_metadatas: list[dict] = []

        for (chunk_id, chunk_content, chunk_index), vector in zip(rows, vectors):
            chroma_ids.append(str(chunk_id))
            chroma_embeddings.append(vector)
            chroma_documents.append(chunk_content)
            chroma_metadatas.append({
                "lesson_id": str(lesson_id),
                "chunk_index": chunk_index,
                "slide_number": str(image_slide_numbers[chunk_index - chunk_offset]),
                "source": "image",
            })

        await asyncio.to_thread(
            lambda: col.add(
                ids=chroma_ids,
                embeddings=chroma_embeddings,
                documents=chroma_documents,
                metadatas=chroma_metadatas,
            )
        )
        await db.commit()

    logger.info("lesson_id=%s ingested: %d image chunk(s)", lesson_id, len(rows))
    return len(rows)


async def search(
    query: str,
    db: AsyncSession,
    k: int = 5,
    lesson_id: int | None = None,
) -> list[str]:
    """
    Embed *query* and return the top-k most similar lesson chunk texts.
    Optionally filter by lesson_id.
    """
    vector = await embed(query)
    col = lesson_chunks_col()

    kwargs: dict = {"query_embeddings": [vector], "n_results": k, "include": ["documents"]}
    if lesson_id is not None:
        kwargs["where"] = {"lesson_id": str(lesson_id)}

    results = col.query(**kwargs)
    docs = results.get("documents") or [[]]
    return list(docs[0])


def delete_lesson_chunks(lesson_id: int) -> None:
    """Remove all ChromaDB entries for *lesson_id*. Call before deleting SQLite rows."""
    col = lesson_chunks_col()
    try:
        col.delete(where={"lesson_id": str(lesson_id)})
    except Exception:
        pass
